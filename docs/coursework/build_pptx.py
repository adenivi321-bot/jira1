"""Builds the defense presentation (.pptx) for the TaskHub coursework.

Output: docs/coursework/Презентация_TaskHub.pptx
Slides: 14 (within the requested 10-15 range).

Style: clean, minimal, bilingual-safe (Russian text, Calibri).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Emu, Pt

HERE = Path(__file__).resolve().parent
PNG = HERE / "diagrams" / "png"
OUT = HERE / "Презентация_TaskHub.pptx"

# 16:9 widescreen
SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)

NAVY = RGBColor(0x0F, 0x2D, 0x52)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
GREY = RGBColor(0x55, 0x5C, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1F, 0x2B)

FONT = "Calibri"


def _txt(tf, text, *, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font=FONT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p


def _add_bullets(tf, items, *, size=18, color=DARK):
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
        p.space_after = Pt(6)


def _bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def _header(slide, title, subtitle=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Cm(1.6))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY

    tb = slide.shapes.add_textbox(Cm(0.8), Cm(0.15), SLIDE_W - Cm(1.6), Cm(1.4))
    tf = tb.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    _txt(tf, title, size=22, bold=True, color=WHITE)

    # subtle accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Cm(1.6), SLIDE_W, Cm(0.08))
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT

    if subtitle:
        sb = slide.shapes.add_textbox(Cm(0.8), Cm(1.85), SLIDE_W - Cm(1.6), Cm(0.9))
        _txt(sb.text_frame, subtitle, size=14, color=GREY)


def _footer(slide, page, total):
    fb = slide.shapes.add_textbox(Cm(0.8), SLIDE_H - Cm(0.9), SLIDE_W - Cm(1.6), Cm(0.7))
    tf = fb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"TaskHub · защита курсовой работы"
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.color.rgb = GREY

    nb = slide.shapes.add_textbox(SLIDE_W - Cm(2.5), SLIDE_H - Cm(0.9), Cm(2.0), Cm(0.7))
    tf2 = nb.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page} / {total}"
    r2.font.name = FONT
    r2.font.size = Pt(10)
    r2.font.color.rgb = GREY


def _content_area(slide):
    # consistent inner content area below header
    return Cm(0.9), Cm(2.9), SLIDE_W - Cm(1.8), SLIDE_H - Cm(4.0)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_divider(prs, page, total, *, kicker, title, summary):
    """Section-divider slide announcing a chapter."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)

    # accent stripe
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(7.6), Cm(3.5), Cm(0.18))
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT

    # kicker
    kb = s.shapes.add_textbox(Cm(2), Cm(6.0), SLIDE_W - Cm(4), Cm(1.4))
    _txt(kb.text_frame, kicker, size=22, bold=True, color=ACCENT)

    # title
    tb = s.shapes.add_textbox(Cm(2), Cm(8.2), SLIDE_W - Cm(4), Cm(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # summary
    sb = s.shapes.add_textbox(Cm(2), Cm(12.0), SLIDE_W - Cm(4), Cm(3.0))
    tf2 = sb.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    r = p.add_run()
    r.text = summary
    r.font.name = FONT
    r.font.size = Pt(18)
    r.font.color.rgb = LIGHT

    # footer with light color on navy bg
    fb = s.shapes.add_textbox(Cm(0.8), SLIDE_H - Cm(0.9), SLIDE_W - Cm(1.6), Cm(0.7))
    p = fb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "TaskHub · защита курсовой работы"
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.color.rgb = LIGHT

    nb = s.shapes.add_textbox(SLIDE_W - Cm(2.5), SLIDE_H - Cm(0.9), Cm(2.0), Cm(0.7))
    p2 = nb.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page} / {total}"
    r2.font.name = FONT
    r2.font.size = Pt(10)
    r2.font.color.rgb = LIGHT


def slide_concept(prs, page, total):
    """1.1 — concept of task management systems."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "1.1. Системы управления задачами",
            "Понятие, назначение, классификация")
    x, y, w, h = _content_area(s)

    # left: definition + functions
    tb = s.shapes.add_textbox(x, y, w * 0.55, h)
    tf = tb.text_frame
    _txt(tf, "Определение", size=18, bold=True, color=NAVY)
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    r = p.add_run()
    r.text = ("Система управления задачами (Issue / Task Tracking System) — программное "
              "средство для регистрации, декомпозиции, планирования и контроля задач "
              "проектной команды.")
    r.font.name = FONT
    r.font.size = Pt(15)
    r.font.color.rgb = DARK

    p2 = tf.add_paragraph()
    p2.space_before = Pt(14)
    r = p2.add_run()
    r.text = "Базовые функции"
    r.font.name = FONT
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = NAVY
    for item in [
        "регистрация задач, ошибок и пользовательских историй;",
        "управление статусами и жизненным циклом задач;",
        "приоритизация и оценка трудозатрат;",
        "коммуникация (комментарии, уведомления, упоминания);",
        "отчётность и метрики (burndown, velocity, lead time).",
    ]:
        p = tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = FONT
        r.font.size = Pt(14)
        r.font.color.rgb = DARK

    # right: classification card
    card_x = x + w * 0.6
    card_w = w * 0.4
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, y, card_w, h * 0.92)
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.adjustments[0] = 0.05

    ct = s.shapes.add_textbox(card_x + Cm(0.5), y + Cm(0.4),
                              card_w - Cm(1.0), h * 0.92 - Cm(0.8))
    tf2 = ct.text_frame
    tf2.word_wrap = True
    _txt(tf2, "Классификация", size=18, bold=True, color=ACCENT)
    classes = [
        ("По модели", "SaaS · on-premise · self-hosted"),
        ("По методологии", "Kanban · Scrum · Waterfall"),
        ("По арендности", "Single-tenant · Multi-tenant"),
        ("Монетизация", "Free · Freemium · Subscription"),
    ]
    for k, v in classes:
        p = tf2.add_paragraph()
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = k + ":  "
        r.font.name = FONT
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = NAVY
        r2 = p.add_run()
        r2.text = v
        r2.font.name = FONT
        r2.font.size = Pt(14)
        r2.font.color.rgb = DARK

    _footer(s, page, total)


def slide_uml_overview(prs, page, total):
    """1.3 — UML 2.5 as a design tool."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "1.3. UML 2.5 как инструмент проектирования",
            "Унифицированный язык моделирования: 14 типов диаграмм")
    x, y, w, h = _content_area(s)

    # 2 columns: structural / behavioral
    col_w = (w - Cm(1.0)) / 2
    cols = [
        ("Структурные диаграммы", NAVY, [
            "Class — классы и связи",
            "Component — компоненты и интерфейсы",
            "Deployment — узлы и артефакты",
            "Package — модульная декомпозиция",
            "Object, Composite, Profile",
        ]),
        ("Поведенческие диаграммы", ACCENT, [
            "Use Case — сценарии использования",
            "Sequence — взаимодействие во времени",
            "Activity — поток управления",
            "State machine — конечный автомат",
            "Communication, Interaction overview, Timing",
        ]),
    ]
    for i, (title, color, items) in enumerate(cols):
        cx = x + i * (col_w + Cm(0.5))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, col_w, h * 0.6)
        card.line.fill.background()
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.adjustments[0] = 0.05

        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, col_w, Cm(0.7))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = color

        tt = s.shapes.add_textbox(cx + Cm(0.3), y + Cm(0.1), col_w - Cm(0.6), Cm(0.6))
        _txt(tt.text_frame, title, size=16, bold=True, color=WHITE)

        lt = s.shapes.add_textbox(cx + Cm(0.4), y + Cm(1.1),
                                  col_w - Cm(0.8), h * 0.6 - Cm(1.4))
        tf = lt.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = "•  " + item
            r.font.name = FONT
            r.font.size = Pt(14)
            r.font.color.rgb = DARK
            p.space_after = Pt(3)

    # bottom callout
    cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            x, y + h * 0.65, w, h * 0.3)
    cb.line.fill.background()
    cb.fill.solid()
    cb.fill.fore_color.rgb = NAVY
    cb.adjustments[0] = 0.07
    ct = s.shapes.add_textbox(x + Cm(0.5), y + h * 0.65 + Cm(0.3),
                              w - Cm(1.0), h * 0.3 - Cm(0.6))
    tf = ct.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Для прототипа TaskHub выбраны: "
    r.font.name = FONT
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    r2 = p.add_run()
    r2.text = ("use case, class, ER, sequence (3 шт.), activity, state, "
               "component, deployment, package — 11 диаграмм, покрывающих структуру и поведение системы.")
    r2.font.name = FONT
    r2.font.size = Pt(15)
    r2.font.color.rgb = WHITE

    _footer(s, page, total)


def slide_requirements(prs, page, total):
    """2.1 — functional requirements + actors via Use Case."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "2.1. Требования к прототипу",
            "Функциональные требования и акторы (UML Use Case)")
    x, y, w, h = _content_area(s)

    # left: requirements list
    tb = s.shapes.add_textbox(x, y, w * 0.42, h)
    tf = tb.text_frame
    _txt(tf, "Функциональные требования", size=16, bold=True, color=NAVY)
    reqs = [
        "регистрация и аутентификация (JWT);",
        "управление арендаторами (tenants);",
        "проекты, доски (Kanban / Scrum), спринты;",
        "CRUD задач, комментарии, вложения;",
        "drag-and-drop по колонкам и спринтам;",
        "история изменений (audit log);",
        "полнотекстовый поиск;",
        "уведомления и e-mail-рассылка;",
        "REST API + WebSocket-обновления.",
    ]
    for item in reqs:
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = FONT
        r.font.size = Pt(13)
        r.font.color.rgb = DARK

    # right: use case diagram
    s.shapes.add_picture(str(PNG / "use-case.png"),
                         x + w * 0.45, y, width=w * 0.55, height=h)

    _footer(s, page, total)


def slide_dynamics(prs, page, total):
    """2.4 — dynamic behavior: sequence + state side by side."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "2.4. Поведение системы",
            "Sequence — создание задачи · State — жизненный цикл")
    x, y, w, h = _content_area(s)

    half_w = (w - Cm(0.5)) / 2
    s.shapes.add_picture(str(PNG / "sequence-create-issue.png"),
                         x, y, width=half_w, height=h * 0.85)
    s.shapes.add_picture(str(PNG / "state-issue.png"),
                         x + half_w + Cm(0.5), y, width=half_w, height=h * 0.85)

    # captions
    c1 = s.shapes.add_textbox(x, y + h * 0.86, half_w, Cm(1.0))
    _txt(c1.text_frame,
         "Sequence: HTTP-путь + асинхронные события (EventEmitter, WebSocket).",
         size=12, color=GREY, align=PP_ALIGN.CENTER)
    c2 = s.shapes.add_textbox(x + half_w + Cm(0.5), y + h * 0.86, half_w, Cm(1.0))
    _txt(c2.text_frame,
         "State: 6 состояний IssueStatus, переходы фиксируются в issue_changelog.",
         size=12, color=GREY, align=PP_ALIGN.CENTER)

    _footer(s, page, total)


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(s, LIGHT)

    # top navy band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Cm(7.5))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY

    # accent diagonal stripe
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Cm(7.5), SLIDE_W, Cm(0.25))
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT

    # institution
    tb = s.shapes.add_textbox(Cm(1.5), Cm(0.8), SLIDE_W - Cm(3), Cm(1.5))
    _txt(tb.text_frame,
         "МИНОБРНАУКИ РОССИИ · Высшее учебное заведение",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # topic label
    tb2 = s.shapes.add_textbox(Cm(1.5), Cm(2.0), SLIDE_W - Cm(3), Cm(1.0))
    _txt(tb2.text_frame, "Защита курсовой работы",
         size=16, color=ACCENT, align=PP_ALIGN.CENTER)

    # title
    tb3 = s.shapes.add_textbox(Cm(1.5), Cm(3.2), SLIDE_W - Cm(3), Cm(3.5))
    tf = tb3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ("Разработка прототипа системы\n"
              "управления задачами (наподобие Jira)\n"
              "с полным комплектом UML-диаграмм")
    r.font.name = FONT
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # author block
    tb4 = s.shapes.add_textbox(Cm(2), Cm(9.5), SLIDE_W - Cm(4), Cm(6))
    tf = tb4.text_frame
    tf.word_wrap = True
    rows = [
        ("Тема:", "Прототип TaskHub — самостоятельно развёртываемая SaaS-система управления задачами"),
        ("Студент:", "____________________________________"),
        ("Группа:", "____________________________________"),
        ("Научный руководитель:", "____________________________________"),
        ("Год:", "2025"),
    ]
    for i, (k, v) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = k + "  "
        r1.font.name = FONT
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = NAVY
        r2 = p.add_run()
        r2.text = v
        r2.font.name = FONT
        r2.font.size = Pt(16)
        r2.font.color.rgb = DARK
        p.space_after = Pt(4)


def slide_relevance(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Актуальность темы")
    x, y, w, h = _content_area(s)

    # left: bullets
    tb = s.shapes.add_textbox(x, y, w * 0.55, h)
    _add_bullets(tb.text_frame, [
        "Распределённые команды требуют единой среды для планирования и контроля задач.",
        "Зарубежные SaaS (Jira, Asana, Linear) — vendor lock-in, рост стоимости лицензий, требования к локализации данных и платежам.",
        "B2B-SaaS как продуктовая модель: multi-tenant изоляция + тарифные планы + интеграция с локальной платёжной системой.",
        "Современный стек (NestJS, React, PostgreSQL) делает разработку такого прототипа реалистичной задачей курсовой работы.",
    ], size=18)

    # right: stats card
    card_x = x + w * 0.6
    card_w = w * 0.4
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, y, card_w, h * 0.85)
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.adjustments[0] = 0.05

    ct = s.shapes.add_textbox(card_x + Cm(0.6), y + Cm(0.4), card_w - Cm(1.2), h - Cm(1.0))
    tf = ct.text_frame
    tf.word_wrap = True
    items = [
        ("85%+", "команд используют ≥1 системы трекинга"),
        ("11", "типов UML-диаграмм в стандарте 2.5"),
        ("4 тарифа", "FREE · BASIC · PRO · ENTERPRISE"),
        ("RLS", "механизм изоляции арендаторов в PostgreSQL"),
    ]
    for i, (big, small) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = big
        r.font.name = FONT
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = ACCENT

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = small
        r2.font.name = FONT
        r2.font.size = Pt(13)
        r2.font.color.rgb = GREY

    _footer(s, page, total)


def slide_goal(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Цель и задачи работы")
    x, y, w, h = _content_area(s)

    # goal
    g = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Cm(2.6))
    g.line.fill.background()
    g.fill.solid()
    g.fill.fore_color.rgb = NAVY
    g.adjustments[0] = 0.08
    gt = s.shapes.add_textbox(x + Cm(0.5), y + Cm(0.3), w - Cm(1.0), Cm(2.0))
    tf = gt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Цель: "
    r.font.name = FONT
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    r2 = p.add_run()
    r2.text = ("разработать прототип многопользовательской системы управления задачами "
               "и подготовить полный комплект UML-диаграмм, описывающих её архитектуру и поведение.")
    r2.font.name = FONT
    r2.font.size = Pt(20)
    r2.font.color.rgb = WHITE

    # tasks
    tt = s.shapes.add_textbox(x, y + Cm(3.0), w, h - Cm(3.0))
    tf2 = tt.text_frame
    _txt(tf2, "Задачи:", size=18, bold=True, color=NAVY)
    tasks = [
        "1. Проанализировать предметную область и существующие решения (Jira, Asana, Linear, Trello).",
        "2. Спроектировать архитектуру: модули, БД, протоколы взаимодействия, изоляцию арендаторов.",
        "3. Реализовать прототип на стеке NestJS + React + PostgreSQL + Redis + MinIO.",
        "4. Построить 11 UML-диаграмм (use case, class, ER, sequence, activity, state, component, deployment, package).",
        "5. Оформить пояснительную записку по требованиям ГОСТ.",
    ]
    for t in tasks:
        p = tf2.add_paragraph()
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(16)
        r.font.color.rgb = DARK
        p.space_after = Pt(4)

    _footer(s, page, total)


def slide_object_methods(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Объект, предмет и методы")
    x, y, w, h = _content_area(s)

    # 3 columns
    col_w = (w - Cm(1.0)) / 3
    cols = [
        ("Объект", "Процессы планирования и контроля задач в распределённых проектных командах.", ACCENT),
        ("Предмет", "Архитектура, модель данных и UML-документация прототипа task-management системы.", NAVY),
        ("Методы", "Системный анализ, объектно-ориентированное проектирование (UML 2.5), реляционное моделирование, программная инженерия.", GREY),
    ]
    for i, (title, body, color) in enumerate(cols):
        cx = x + i * (col_w + Cm(0.5))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, col_w, h * 0.9)
        card.line.fill.background()
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.adjustments[0] = 0.05

        # color stripe at top of card
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, col_w, Cm(0.3))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = color

        tb = s.shapes.add_textbox(cx + Cm(0.4), y + Cm(0.6), col_w - Cm(0.8), h * 0.9 - Cm(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = FONT
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.space_before = Pt(12)
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = FONT
        r2.font.size = Pt(16)
        r2.font.color.rgb = DARK

    _footer(s, page, total)


def slide_compare(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Сравнение существующих решений", "Анализ аналогов задал требования к прототипу TaskHub")
    x, y, w, h = _content_area(s)

    headers = ["Система", "Поставка", "Multi-tenant", "Тарифы", "Гибкость WF"]
    rows = [
        ["Atlassian Jira",  "Cloud + DC",   "Нет",      "Per-user",          "Высокая"],
        ["GitHub Projects", "SaaS",         "Нет",      "С GitHub",          "Средняя"],
        ["Asana / Trello",  "SaaS",         "Нет",      "Free / Premium",    "Низкая"],
        ["YouTrack",        "Cloud + on-prem", "Огранич.", "Per-user",       "Высокая"],
        ["OpenProject",     "On-prem",      "Нет",      "Community / Ent.",  "Средняя"],
        ["TaskHub (proto)", "On-prem (Docker)", "Да (RLS)", "4 плана + YooKassa", "Высокая"],
    ]

    rows_n = len(rows) + 1
    cols_n = len(headers)
    table_w = w
    table_h = h * 0.85
    tbl_shape = s.shapes.add_table(rows_n, cols_n, x, y + Cm(0.4), table_w, table_h)
    tbl = tbl_shape.table

    col_widths = [0.28, 0.18, 0.18, 0.18, 0.18]
    for i, frac in enumerate(col_widths):
        tbl.columns[i].width = int(table_w * frac)

    for j, header in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Cm(0.2)
        tf.margin_right = Cm(0.2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = header
        r.font.name = FONT
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = WHITE

    for i, row in enumerate(rows, start=1):
        is_taskhub = row[0].startswith("TaskHub")
        for j, value in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            if is_taskhub:
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF1, 0xFA)
            else:
                cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Cm(0.2)
            tf.margin_right = Cm(0.2)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = value
            r.font.name = FONT
            r.font.size = Pt(14)
            r.font.bold = is_taskhub or j == 0
            r.font.color.rgb = NAVY if is_taskhub else DARK

    _footer(s, page, total)


def slide_diagram(prs, page, total, *, title, subtitle, png, bullets):
    """Generic slide with diagram on the left and bullets on the right."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, title, subtitle)
    x, y, w, h = _content_area(s)

    # diagram occupies ~62% width
    diag_w = w * 0.6
    diag_h = h
    s.shapes.add_picture(str(png), x, y, width=diag_w, height=diag_h)

    # bullets on right
    tb = s.shapes.add_textbox(x + diag_w + Cm(0.6), y + Cm(0.2),
                              w - diag_w - Cm(0.6), h - Cm(0.4))
    _add_bullets(tb.text_frame, bullets, size=15)

    _footer(s, page, total)


def slide_architecture_text(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Архитектура прототипа TaskHub",
            "Трёхуровневая SPA-архитектура с асинхронной шиной событий")
    x, y, w, h = _content_area(s)

    # left: bullets
    tb = s.shapes.add_textbox(x, y, w * 0.5, h)
    _add_bullets(tb.text_frame, [
        "Frontend: React 18 + Vite + TanStack Query + Zustand.",
        "Backend: NestJS + TypeORM, 10 модулей (Auth, Tenants, Users, Projects, Issues, Search, Notifications, Health, Billing, Redis).",
        "Хранилища: PostgreSQL 16 (RLS), Redis (кеш и pub/sub), MinIO (вложения).",
        "Биллинг: YooKassa-интеграция + PlanGuard, ограничения по тарифам (projects/users/issues).",
        "Реал-тайм: WebSocket (Socket.IO) + EventEmitter для развязки слоёв.",
        "Развёртывание: Docker Compose, обратный прокси Traefik.",
    ], size=15)

    # right: component diagram
    s.shapes.add_picture(str(PNG / "component.png"),
                         x + w * 0.55, y, width=w * 0.45, height=h)

    _footer(s, page, total)


def slide_db(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Модель данных и многоарендность",
            "ER-схема и изоляция через Row-Level Security")
    x, y, w, h = _content_area(s)

    s.shapes.add_picture(str(PNG / "er.png"), x, y, width=w * 0.62, height=h)
    tb = s.shapes.add_textbox(x + w * 0.65, y + Cm(0.2),
                              w - w * 0.65, h - Cm(0.4))
    _add_bullets(tb.text_frame, [
        "12 таблиц, все доменные сущности привязаны к tenants(id).",
        "tenant_id обязателен и индексирован — без него любой запрос отклоняется RLS-политикой.",
        "UUID-первичные ключи, FK с ON DELETE CASCADE для целостности.",
        "Полнотекстовый поиск по issues через tsvector + GIN-индекс.",
        "Fractional indexing (column_position) для drag-and-drop без массовых апдейтов.",
    ], size=15)

    _footer(s, page, total)


def slide_tech(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Технологический стек",
            "Обоснованный выбор зрелых open-source инструментов")
    x, y, w, h = _content_area(s)

    groups = [
        ("Frontend", ACCENT, [
            "React 18", "TypeScript", "Vite",
            "TanStack Query", "Zustand", "Tailwind CSS",
        ]),
        ("Backend", NAVY, [
            "NestJS 10", "TypeORM", "class-validator",
            "Passport JWT", "Socket.IO", "BullMQ",
        ]),
        ("Инфраструктура", GREY, [
            "PostgreSQL 16", "Redis 7", "MinIO",
            "Docker Compose", "Traefik", "Prometheus",
        ]),
    ]
    col_w = (w - Cm(1.0)) / 3
    for i, (name, color, items) in enumerate(groups):
        cx = x + i * (col_w + Cm(0.5))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, col_w, h * 0.9)
        card.line.fill.background()
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.adjustments[0] = 0.05

        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, col_w, Cm(0.7))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = color

        title_tb = s.shapes.add_textbox(cx + Cm(0.3), y + Cm(0.1), col_w - Cm(0.6), Cm(0.6))
        _txt(title_tb.text_frame, name, size=16, bold=True, color=WHITE)

        list_tb = s.shapes.add_textbox(cx + Cm(0.4), y + Cm(1.1),
                                       col_w - Cm(0.8), h * 0.9 - Cm(1.4))
        tf = list_tb.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = "•  " + item
            r.font.name = FONT
            r.font.size = Pt(16)
            r.font.color.rgb = DARK
            p.space_after = Pt(4)

    _footer(s, page, total)


def slide_results(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Результаты работы")
    x, y, w, h = _content_area(s)

    metrics = [
        ("11", "UML-диаграмм в полном комплекте"),
        ("9", "NestJS-модулей backend"),
        ("12", "таблиц БД с RLS-изоляцией"),
        ("4", "роли пользователей в системе"),
    ]
    col_w = (w - Cm(2.0)) / 4
    for i, (big, small) in enumerate(metrics):
        cx = x + i * (col_w + Cm(0.5))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, col_w, Cm(4.5))
        card.line.fill.background()
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY
        card.adjustments[0] = 0.07

        big_tb = s.shapes.add_textbox(cx, y + Cm(0.5), col_w, Cm(2.5))
        _txt(big_tb.text_frame, big, size=54, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
        small_tb = s.shapes.add_textbox(cx + Cm(0.3), y + Cm(3.0),
                                        col_w - Cm(0.6), Cm(1.4))
        _txt(small_tb.text_frame, small, size=13, color=WHITE,
             align=PP_ALIGN.CENTER)

    # bullets below
    tb = s.shapes.add_textbox(x, y + Cm(5.2), w, h - Cm(5.4))
    _add_bullets(tb.text_frame, [
        "Реализован прототип системы: аутентификация, проекты, доски, задачи, комментарии, вложения, поиск.",
        "Подготовлена пояснительная записка (≈30 страниц) по требованиям ГОСТ 7.32.",
        "Архитектура верифицирована UML-моделями уровня use case, classes, ER, dynamic behavior, deployment.",
        "Стек контейнеризован: проект разворачивается единой командой docker compose up.",
    ], size=15)

    _footer(s, page, total)


def slide_conclusion(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Заключение и перспективы")
    x, y, w, h = _content_area(s)

    # left: outcomes
    tb1 = s.shapes.add_textbox(x, y, w * 0.5, h)
    tf = tb1.text_frame
    _txt(tf, "Достигнуто:", size=20, bold=True, color=NAVY)
    for item in [
        "цель и задачи курсовой работы выполнены в полном объёме;",
        "построен прототип, демонстрирующий ключевые сценарии Jira-подобной системы;",
        "комплект UML-диаграмм покрывает структуру и поведение системы;",
        "изоляция арендаторов реализована на уровне БД через RLS.",
    ]:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = "✓  " + item
        r.font.name = FONT
        r.font.size = Pt(15)
        r.font.color.rgb = DARK

    # right: next steps
    tb2 = s.shapes.add_textbox(x + w * 0.55, y, w * 0.45, h)
    tf2 = tb2.text_frame
    _txt(tf2, "Перспективы развития:", size=20, bold=True, color=ACCENT)
    for item in [
        "интеграции с Git-провайдерами (GitHub, GitLab) и CI;",
        "автоматизации рабочих процессов (правила, триггеры);",
        "отчёты, диаграммы выгорания, прогноз сроков;",
        "мобильный клиент и offline-режим;",
        "масштабирование на Kubernetes и horizontal sharding по tenant_id.",
    ]:
        p = tf2.add_paragraph()
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = "→  " + item
        r.font.name = FONT
        r.font.size = Pt(15)
        r.font.color.rgb = DARK

    _footer(s, page, total)


def slide_thanks(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)

    tb = s.shapes.add_textbox(Cm(2), Cm(6.5), SLIDE_W - Cm(4), Cm(3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Спасибо за внимание!"
    r.font.name = FONT
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = WHITE

    tb2 = s.shapes.add_textbox(Cm(2), Cm(11.0), SLIDE_W - Cm(4), Cm(2))
    tf2 = tb2.text_frame
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Готов ответить на вопросы"
    r.font.name = FONT
    r.font.size = Pt(22)
    r.font.color.rgb = ACCENT

    _footer(s, page, total)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 15 slides total (within the requested 10-15 range).
    # Structure mirrors the coursework: Введение → Глава 1 (теоретическая
    # часть) → Глава 2 (практическая часть) → Заключение.
    TOTAL = 15

    # — Введение —
    slide_title(prs)                                           # 1
    slide_relevance(prs, 2, TOTAL)                             # 2
    slide_goal(prs, 3, TOTAL)                                  # 3

    # — Глава 1. Теоретическая часть —
    slide_divider(prs, 4, TOTAL,                               # 4
                  kicker="ГЛАВА 1",
                  title="Теоретическая часть",
                  summary=("Понятие систем управления задачами, анализ "
                           "аналогов, обоснование выбора UML 2.5 и "
                           "технологического стека прототипа."))
    slide_concept(prs, 5, TOTAL)                               # 5  · 1.1
    slide_compare(prs, 6, TOTAL)                               # 6  · 1.2
    slide_uml_overview(prs, 7, TOTAL)                          # 7  · 1.3
    slide_tech(prs, 8, TOTAL)                                  # 8  · 1.4

    # — Глава 2. Практическая часть —
    slide_divider(prs, 9, TOTAL,                               # 9
                  kicker="ГЛАВА 2",
                  title="Практическая часть",
                  summary=("Разработка прототипа TaskHub: требования, "
                           "архитектура, модель данных, поведение системы "
                           "и развёртывание."))
    slide_requirements(prs, 10, TOTAL)                         # 10 · 2.1 (Use Case)
    slide_architecture_text(prs, 11, TOTAL)                    # 11 · 2.2 (Component)
    slide_diagram(prs, 12, TOTAL,                              # 12 · 2.3 (Class)
                  title="2.3. Доменная модель (UML Class)",
                  subtitle="11 сущностей и 5 перечислений TaskHub",
                  png=PNG / "class.png",
                  bullets=[
                      "Tenant, User, Project, Sprint, Issue, Comment, Attachment …",
                      "Перечисления: UserRole, IssueStatus, IssuePriority, BoardType.",
                      "Композиция Tenant → Project → Issue фиксирует владение.",
                      "Соответствует TypeORM-сущностям backend-модулей.",
                  ])
    slide_db(prs, 13, TOTAL)                                   # 13 · 2.4 (ER + RLS)
    slide_dynamics(prs, 14, TOTAL)                             # 14 · 2.5 (Sequence + State)

    # — Заключение —
    slide_conclusion(prs, 15, TOTAL)                           # 15

    prs.save(OUT)
    print(f"Saved: {OUT}  ({OUT.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    build()
